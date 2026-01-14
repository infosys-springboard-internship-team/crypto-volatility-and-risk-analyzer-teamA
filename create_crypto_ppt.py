from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation
prs = Presentation()

# Define dark background slide layout
def add_slide(title, content):
    slide_layout = prs.slide_layouts[6]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 15, 36)  # Dark blue background

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(88, 166, 255)  # Light blue

    # Content
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
    tf2 = txBox2.text_frame
    for line in content.split("\n"):
        p2 = tf2.add_paragraph()
        p2.text = line
        p2.font.size = Pt(22)
        p2.font.color.rgb = RGBColor(255, 255, 255)

# Slide 1: Title Slide
add_slide(
    "Crypto Volatility & Risk Analyzer – Milestone 1",
    "Presented by: A. Hemalatha\nInfosys Springboard Internship"
)

# Slide 2: Project Overview
add_slide(
    "Project Overview",
    "- Analyze cryptocurrency volatility and risk\n- Fetch live and historical crypto data\n- Calculate returns, volatility, and Value at Risk (VaR)\n- Visualize trends, returns, and trading volumes"
)

# Slide 3: Data Acquisition
add_slide(
    "Data Acquisition",
    "- Used yFinance for Bitcoin historical data (1 year)\n- Used CoinGecko API for live prices of 8 cryptocurrencies\n- Stored data in CSV files for analysis"
)

# Slide 4: Daily Returns & Volatility
add_slide(
    "Daily Returns & Volatility",
    "- Calculated daily returns using percentage change\n- Daily volatility: standard deviation of returns\n- Annual volatility: scaled by sqrt(252)"
)

# Slide 5: Value at Risk (VaR)
add_slide(
    "Value at Risk (VaR)",
    "- 5% confidence level\n- Measures potential loss in worst-case scenario\n- Calculated using historical return quantiles"
)

# Slide 6: Visualizations
add_slide(
    "Visualizations",
    "- Price trend over 1 year\n- Return distribution\n- Cumulative returns\n- 7-day trends for BTC, ETH, SOL\n- 24h trading volume comparison"
)

# Slide 7: Live Crypto Dashboard
add_slide(
    "Live Crypto Dashboard",
    "- Built using Streamlit and Plotly\n- Real-time price, 24h change, and volume\n- Trend charts for selected cryptocurrencies\n- Interactive and refreshable UI"
)

# Slide 8: Conclusion
add_slide(
    "Conclusion",
    "- Successfully implemented Milestone 1\n- Fetched, processed, and visualized crypto data\n- Ready for analysis of volatility and risk in next milestones"
)

# Slide 9: Thank You
add_slide(
    "Thank You",
    "Prepared by: A. Hemalatha\nInfosys Springboard Internship"
)

# Save the presentation
prs.save("Crypto_Volatility_Risk_Analyzer.pptx")
print("Presentation created successfully!")
